#!/usr/bin/env python3
"""
extract_text.py - 从 Office 文档和 PDF 中提取纯文本内容。

支持格式：.docx, .xlsx, .pptx, .pdf
依赖库：按需安装，缺少时给出提示并优雅退出。

用法：
    python3 extract_text.py --input "path/to/file.docx" --output "path/to/output.txt"
"""

import argparse
import sys
import os


def extract_docx(filepath):
    """从 .docx 文件提取文本"""
    try:
        from docx import Document
    except ImportError:
        print("ERROR: 缺少 python-docx 库。请运行: pip install python-docx", file=sys.stderr)
        sys.exit(1)

    doc = Document(filepath)
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    return "\n".join(paragraphs)


def extract_xlsx(filepath):
    """从 .xlsx 文件提取文本"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        print("ERROR: 缺少 openpyxl 库。请运行: pip install openpyxl", file=sys.stderr)
        sys.exit(1)

    wb = load_workbook(filepath, read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip(" |"):
                lines.append(row_text)
        lines.append("")
    wb.close()
    return "\n".join(lines)


def extract_pptx(filepath):
    """从 .pptx 文件提取文本"""
    try:
        from pptx import Presentation
    except ImportError:
        print("ERROR: 缺少 python-pptx 库。请运行: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(filepath)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        lines.append(row_text)
        lines.append("")
    return "\n".join(lines)


def extract_pdf(filepath):
    """从 .pdf 文件提取文本"""
    try:
        import pdfplumber
    except ImportError:
        try:
            from PyPDF2 import PdfReader
        except ImportError:
            print("ERROR: 缺少 PDF 解析库。请运行: pip install pdfplumber 或 pip install PyPDF2", file=sys.stderr)
            sys.exit(1)

        # PyPDF2 fallback
        reader = PdfReader(filepath)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                pages.append(text.strip())
        return "\n\n".join(pages)

    # pdfplumber (preferred)
    lines = []
    with pdfplumber.open(filepath) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                lines.append(text.strip())
            # 也提取表格
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    row_text = " | ".join(cell if cell else "" for cell in row)
                    if row_text.strip(" |"):
                        lines.append(row_text)
    return "\n\n".join(lines)


EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
}


def main():
    parser = argparse.ArgumentParser(description="从 Office/PDF 文档提取纯文本")
    parser.add_argument("--input", required=True, help="输入文件路径")
    parser.add_argument("--output", required=True, help="输出文本文件路径")
    args = parser.parse_args()

    input_path = args.input
    output_path = args.output

    if not os.path.exists(input_path):
        print(f"ERROR: 文件不存在: {input_path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(input_path)[1].lower()
    if ext not in EXTRACTORS:
        print(f"ERROR: 不支持的文件格式: {ext}。支持: {', '.join(EXTRACTORS.keys())}", file=sys.stderr)
        sys.exit(1)

    try:
        text = EXTRACTORS[ext](input_path)
    except Exception as e:
        print(f"ERROR: 提取文本失败: {e}", file=sys.stderr)
        sys.exit(1)

    if not text or not text.strip():
        print(f"WARNING: 文件 {input_path} 未提取到有效文本内容", file=sys.stderr)
        text = ""

    # 确保输出目录存在
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

    print(f"OK: 已提取 {len(text)} 字符到 {output_path}")


if __name__ == "__main__":
    main()
